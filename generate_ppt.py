from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Colors ──────────────────────────────────────────────
BG_DARK   = RGBColor(0x0A, 0x0E, 0x1B)   # deep navy
BG_CARD   = RGBColor(0x12, 0x18, 0x2E)   # card bg
ACCENT    = RGBColor(0x00, 0xD4, 0xFF)   # cyan neon
ACCENT2   = RGBColor(0x00, 0x77, 0xFF)   # blue
ACCENT3   = RGBColor(0x7C, 0x3A, 0xED)   # purple
ACCENT4   = RGBColor(0x10, 0xB9, 0x81)   # green
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0x8E, 0x94, 0xA8)
LIGHT_GRAY= RGBColor(0x64, 0x6A, 0x7E)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ── Helpers ─────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None, corner_radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        if border_width:
            shape.line.width = border_width
    if corner_radius:
        shape.adjustments[0] = corner_radius
    return shape

def add_round_rect(slide, left, top, width, height, fill_color=None, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        shape.line.width = Pt(1)
    return shape

def add_textbox(slide, left, top, width, height, text="", font_size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_para(tf, text, size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, space_before=0, space_after=0):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Microsoft YaHei"
    p.alignment = alignment
    if space_before:
        p.space_before = Pt(space_before)
    if space_after:
        p.space_after = Pt(space_after)
    return p

def add_line(slide, x1, y1, x2, y2, color=ACCENT, width=Pt(1)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # MSO_CONNECTOR.STRAIGHT = 1
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_tech_corners(slide, left, top, width, height, color=ACCENT, size=Inches(0.15)):
    """Decorative corner brackets"""
    for x, y, flip_h, flip_v in [
        (left, top, 1, 1),
        (left + width - size, top, -1, 1),
        (left, top + height - size, 1, -1),
        (left + width - size, top + height - size, -1, -1),
    ]:
        # horizontal
        add_rect(slide, x, y, size, Pt(2), fill_color=color)
        # vertical
        add_rect(slide, x, y, Pt(2), size, fill_color=color)

def new_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, BG_DARK)
    return slide

# ── Slide 1: Title ──────────────────────────────────────
slide = new_slide()

# Decorative grid lines
for i in range(14):
    x = Inches(i * 1.0)
    add_line(slide, x, 0, x, H, color=RGBColor(0x15, 0x1D, 0x35), width=Pt(0.5))
for i in range(8):
    y = Inches(i * 1.0)
    add_line(slide, 0, y, W, y, color=RGBColor(0x15, 0x1D, 0x35), width=Pt(0.5))

# Corner brackets
add_tech_corners(slide, Inches(0.8), Inches(0.8), Inches(11.7), Inches(5.8), color=ACCENT, size=Inches(0.3))

# Title
tf = add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
                 "AI 信息周报", font_size=56, color=WHITE, bold=True)
add_para(tf, "Artificial Intelligence Weekly Report", size=22, color=GRAY, space_before=8)

# Accent line
add_line(slide, Inches(1.5), Inches(3.4), Inches(5.5), Inches(3.4), color=ACCENT, width=Pt(3))

# Date & meta
add_textbox(slide, Inches(1.5), Inches(3.7), Inches(6), Inches(0.6),
            "2026/06/08  ·  第 23 周  ·  周一发布", font_size=16, color=LIGHT_GRAY)

# Bottom tags
tags = ["#商业化破局", "#超级入口重构", "#具身智能", "#端侧AI", "#AI安全"]
for i, tag in enumerate(tags):
    x = Inches(1.5 + i * 2.1)
    tag_box = add_round_rect(slide, x, Inches(5.4), Inches(1.9), Inches(0.45), fill_color=BG_CARD, border_color=ACCENT)
    tag_box.line.width = Pt(1)
    tf = tag_box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = tag
    p.font.size = Pt(11)
    p.font.color.rgb = ACCENT
    p.font.name = "SF Mono"
    p.alignment = PP_ALIGN.CENTER

# ── Slide 2: Overview ───────────────────────────────────
slide = new_slide()

# Top accent bar
add_rect(slide, 0, 0, W, Pt(3), fill_color=ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.7),
            "本周速览", font_size=36, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.1), Inches(8), Inches(0.5),
            "WEEKLY OVERVIEW  ·  7 条关键信息", font_size=14, color=GRAY)

add_line(slide, Inches(0.8), Inches(1.7), Inches(3.5), Inches(1.7), color=ACCENT, width=Pt(2))

# Overview cards - 2 rows x 4 cols layout
items = [
    ("01", "豆包正式收费", "月活3.45亿\n三档付费方案", ACCENT),
    ("02", "OpenAI 机器人", "具身智能新部门\n直面特斯拉竞争", ACCENT2),
    ("03", "英伟达 AI PC", "N1X芯片发布\n端侧AI落地", ACCENT3),
    ("04", "宇树科技过会", "73天闪电过会\n人形机器人第一股", ACCENT4),
    ("05", "GPT-5.5 推送", "幻觉率降52.5%\n推理速度×3", ACCENT),
    ("06", "微信 AI Agent", "14亿用户入口\n本月上线测试", ACCENT2),
    ("07", "Claude Code 漏洞", "AI供应链安全\n工具链攻击预警", ACCENT3),
]

card_w = Inches(2.8)
card_h = Inches(2.1)
start_x = Inches(0.8)
start_y = Inches(2.2)
gap_x = Inches(0.15)
gap_y = Inches(0.15)

for i, (num, title, desc, color) in enumerate(items):
    row = i // 4
    col = i % 4
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)

    card = add_round_rect(slide, x, y, card_w, card_h, fill_color=BG_CARD, border_color=color)
    card.line.width = Pt(1)

    # Number
    add_textbox(slide, x + Inches(0.2), y + Inches(0.15), Inches(0.8), Inches(0.5),
                num, font_size=28, color=color, bold=True)

    # Title
    add_textbox(slide, x + Inches(0.2), y + Inches(0.7), card_w - Inches(0.4), Inches(0.5),
                title, font_size=16, color=WHITE, bold=True)

    # Desc
    add_textbox(slide, x + Inches(0.2), y + Inches(1.2), card_w - Inches(0.4), Inches(0.7),
                desc, font_size=11, color=GRAY)

# Bottom note
add_textbox(slide, Inches(0.8), Inches(6.8), Inches(10), Inches(0.4),
            "数据来源：公开报道整理  ·  AI周报团队出品", font_size=10, color=LIGHT_GRAY)

# ── News Detail Slides (3-9) ────────────────────────────
news = [
    {
        "num": "01",
        "title": "豆包 6 月下旬正式收费",
        "tag": "商业化",
        "color": ACCENT,
        "highlight": "月活 3.45 亿的豆包，终于迈出商业化第一步。",
        "cards": [
            ("标准版", "¥68/月", "日常辅助"),
            ("加强版", "¥200/月", "专业创作"),
            ("专业版", "¥500/月", "企业级应用"),
        ],
        "points": [
            "日均 Token 消耗从 1200 亿飙升至 120 万亿，算力成本压力巨大",
            "基础功能继续免费，付费针对 PPT 生成、数据分析、影视制作等复杂任务",
            "腾讯元宝、阿里千问暂未公布收费计划",
        ],
        "date": "2026-06-02",
    },
    {
        "num": "02",
        "title": "OpenAI 成立机器人部门",
        "tag": "具身智能",
        "color": ACCENT2,
        "highlight": "\"生成式 AI → 世界模型 → 机器人执行\"——完整路径浮出水面。",
        "cards": [
            ("负责人", "Aditya Ramesh", "Sora 核心开发者"),
            ("路线", "工业 → 消费级", "两步走战略"),
            ("竞争", "Optimus / Figure", "直接竞争"),
        ],
        "points": [
            "世界模拟能力 × 物理机器人 = 具身智能",
            "由 Sora 核心开发者 Aditya Ramesh 带队",
            "从工业机器人起步，最终造个人消费级机器人",
        ],
        "date": "2026-06-02",
    },
    {
        "num": "03",
        "title": "英伟达发布 RTX Spark AI PC",
        "tag": "端侧AI",
        "color": ACCENT3,
        "highlight": "黄仁勋 GTC 台北：AI 的终点不在云端，而在你桌面上。",
        "cards": [
            ("N1X 芯片", "3nm / 128GB", "联发科联合研发"),
            ("Isaac GR00T", "机器人平台", "人形机器人开发"),
            ("Cosmos 3", "全模态开源", "物理AI基础模型"),
        ],
        "points": [
            "本地直接运行 AI 智能体，无需联网",
            "华硕、戴尔、联想等秋季推出相关产品",
            "全球首款完全开源的全模态物理 AI 基础模型",
        ],
        "date": "2026-06-02",
    },
    {
        "num": "04",
        "title": "宇树科技 73 天闪电过会",
        "tag": "资本市场",
        "color": ACCENT4,
        "highlight": "A 股\"人形机器人第一股\"落定，赛道正式进入资本时代。",
        "cards": [
            ("过会速度", "73 天", "受理→过会"),
            ("募资规模", "42.02 亿", "拟募资金额"),
            ("Q1 营收", "4.23 亿", "2026 第一季度"),
        ],
        "points": [
            "腾讯、阿里持股，产业资本加持",
            "人形机器人赛道正式进入资本市场",
            "标志性事件：从产业爆发到资本闭环",
        ],
        "date": "2026-06-02",
    },
    {
        "num": "05",
        "title": "OpenAI GPT-5.5 全量推送",
        "tag": "模型突破",
        "color": ACCENT,
        "highlight": "AI 正式进入「生产级」时代。",
        "cards": [
            ("幻觉率", "↓ 52.5%", "断崖式下降"),
            ("推理速度", "× 3 倍", "显著提升"),
            ("上下文", "100 万 Token", "突破上限"),
        ],
        "points": [
            "推理速度与 GPT-5.4 持平，但 Token 用量减少",
            "响应更快捷，成本更低",
            "标志着 AI 从\"可用\"迈向\"可靠\"",
        ],
        "date": "2026-05-28",
    },
    {
        "num": "06",
        "title": "微信即将推出 AI Agent",
        "tag": "超级入口",
        "color": ACCENT2,
        "highlight": "从\"连接人与服务\"升级为\"Agent 代替人去调用服务\"。",
        "cards": [
            ("用户基数", "14 亿", "微信全量用户"),
            ("核心能力", "调用小程序", "自动完成任务"),
            ("时间线", "本月上线", "最快6月测试"),
        ],
        "points": [
            "应用内自动调用小程序完成各类复杂任务",
            "超级 App 的流量分配逻辑彻底改写",
            "腾讯借微信生态实现 Agent 弯道超车",
        ],
        "date": "2026-06-02",
    },
    {
        "num": "07",
        "title": "微软警告：Claude Code 存在漏洞",
        "tag": "AI安全",
        "color": ACCENT3,
        "highlight": "AI 供应链安全敲响警钟——攻击面从\"模型\"扩展到\"工具链\"。",
        "cards": [
            ("发现方", "微软威胁情报", "安全团队"),
            ("漏洞类型", "提示词注入", "CI/CD 攻击"),
            ("影响", "凭证泄露", "GitHub Token"),
        ],
        "points": [
            "攻击者可通过提示词注入窃取敏感凭证",
            "AI 辅助开发场景下的 AI 供应链安全典型案例",
            "安全边界从\"模型安全\"扩展到全工具链",
        ],
        "date": "2026-06-07",
    },
]

for item in news:
    slide = new_slide()
    c = item["color"]

    # Top accent bar
    add_rect(slide, 0, 0, W, Pt(3), fill_color=c)

    # Number + Tag
    add_textbox(slide, Inches(0.8), Inches(0.5), Inches(1), Inches(0.5),
                item["num"], font_size=36, color=c, bold=True)
    tag_badge = add_round_rect(slide, Inches(1.8), Inches(0.6), Inches(1.2), Inches(0.4), fill_color=c)
    tf = tag_badge.text_frame
    p = tf.paragraphs[0]
    p.text = item["tag"]
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.CENTER

    # Title
    add_textbox(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.7),
                item["title"], font_size=32, color=WHITE, bold=True)

    # Highlight quote
    quote_box = add_round_rect(slide, Inches(0.8), Inches(2.0), Inches(11.7), Inches(0.7), fill_color=BG_CARD, border_color=c)
    quote_box.line.width = Pt(1)
    tf = quote_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "💡  " + item["highlight"]
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.name = "Microsoft YaHei"
    p.alignment = PP_ALIGN.LEFT

    # Data cards (3 columns)
    card_w2 = Inches(3.7)
    card_h2 = Inches(1.2)
    start_x2 = Inches(0.8)
    gap2 = Inches(0.3)
    for j, (label, value, sub) in enumerate(item["cards"]):
        x = start_x2 + j * (card_w2 + gap2)
        y = Inches(3.1)
        c2 = add_round_rect(slide, x, y, card_w2, card_h2, fill_color=BG_CARD, border_color=c)
        c2.line.width = Pt(1)

        add_textbox(slide, x + Inches(0.2), y + Inches(0.1), card_w2 - Inches(0.4), Inches(0.3),
                    label, font_size=11, color=GRAY)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.4), card_w2 - Inches(0.4), Inches(0.4),
                    value, font_size=24, color=c, bold=True)
        add_textbox(slide, x + Inches(0.2), y + Inches(0.85), card_w2 - Inches(0.4), Inches(0.3),
                    sub, font_size=11, color=LIGHT_GRAY)

    # Key points
    y_points = Inches(4.7)
    add_textbox(slide, Inches(0.8), y_points, Inches(3), Inches(0.4),
                "关键解读", font_size=14, color=c, bold=True)
    for j, point in enumerate(item["points"]):
        # bullet dot
        add_rect(slide, Inches(1.0), y_points + Inches(0.55 + j * 0.45), Pt(6), Pt(6), fill_color=c)
        add_textbox(slide, Inches(1.2), y_points + Inches(0.45 + j * 0.45), Inches(11), Inches(0.45),
                    point, font_size=13, color=WHITE)

    # Date
    add_textbox(slide, Inches(0.8), Inches(6.8), Inches(4), Inches(0.4),
                f"发布时间：{item['date']}", font_size=11, color=LIGHT_GRAY)

    # Page number
    add_textbox(slide, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.4),
                f"{int(item['num'])+1} / 9", font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

# ── Slide 10: Key Takeaways ─────────────────────────────
slide = new_slide()

add_rect(slide, 0, 0, W, Pt(3), fill_color=ACCENT)

add_textbox(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.7),
            "本周要点", font_size=36, color=WHITE, bold=True)
add_textbox(slide, Inches(0.8), Inches(1.1), Inches(8), Inches(0.5),
            "KEY TAKEAWAYS", font_size=14, color=GRAY)

add_line(slide, Inches(0.8), Inches(1.7), Inches(3.5), Inches(1.7), color=ACCENT, width=Pt(2))

takeaways = [
    ("💰", "商业化破局", "从\"烧钱抢流量\"到\"为算力买单\"\n豆包收费是行业风向标", ACCENT),
    ("🔗", "超级入口重构", "Agent 重塑超级 App\n流量分配逻辑彻底改写", ACCENT2),
    ("🤖", "具身智能大爆发", "\"世界模型\"与\"物理躯体\"合体\n软件定义硬件的时代来临", ACCENT3),
    ("💻", "端侧与入口争夺", "Agent 成为新一代操作系统\n从云到端的战场全面铺开", ACCENT4),
    ("🛡️", "安全新战场", "AI 安全从\"模型安全\"扩展到\n\"AI 供应链安全\"", ACCENT),
]

for i, (icon, title, desc, color) in enumerate(takeaways):
    y = Inches(2.2 + i * 1.0)

    # Card
    card = add_round_rect(slide, Inches(0.8), y, Inches(11.7), Inches(0.85), fill_color=BG_CARD, border_color=color)
    card.line.width = Pt(1)

    # Left accent bar
    add_rect(slide, Inches(0.8), y, Pt(4), Inches(0.85), fill_color=color)

    # Icon
    add_textbox(slide, Inches(1.2), y + Inches(0.15), Inches(0.6), Inches(0.5),
                icon, font_size=28, color=WHITE)

    # Title
    add_textbox(slide, Inches(2.0), y + Inches(0.05), Inches(3), Inches(0.4),
                title, font_size=20, color=color, bold=True)

    # Desc
    add_textbox(slide, Inches(2.0), y + Inches(0.4), Inches(9), Inches(0.45),
                desc, font_size=12, color=GRAY)

# Bottom
add_textbox(slide, Inches(0.8), Inches(6.8), Inches(10), Inches(0.4),
            "AI 周报团队  ·  每周一发布  ·  欢迎订阅", font_size=10, color=LIGHT_GRAY)

# ── Save ────────────────────────────────────────────────
output_path = "/Users/chan/AI信息周报/AI周报_PPT_20260608.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
print(f"Slides: {len(prs.slides)}")
